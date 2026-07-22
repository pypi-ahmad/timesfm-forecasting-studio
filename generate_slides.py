# /// script
# requires-python = ">=3.14,<3.15"
# dependencies = ["python-pptx==1.0.2"]
# ///
"""Generate the TimesFM Masterclass presentation from the tutorial Markdown files."""

from __future__ import annotations

import argparse
import math
import os
import re
import tempfile
from dataclasses import dataclass
from pathlib import Path
from urllib.parse import urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
BACKGROUND = "0F172A"
CARD = "1E293B"
CARD_ALT = "172033"
WHITE = "FFFFFF"
BODY = "E2E8F0"
MUTED = "94A3B8"
TEAL = "14B8A6"
TEAL_DARK = "0F766E"
AMBER = "F59E0B"
ROSE = "FB7185"
BLUE = "60A5FA"
FONT = "Calibri"
CODE_FONT = "Courier New"


def rgb(hex_color: str) -> RGBColor:
    """Convert a six-digit hexadecimal color to a python-pptx RGBColor."""
    return RGBColor.from_string(hex_color)


def normalize_heading(value: str) -> str:
    value = re.sub(r"[*_`]", "", value)
    value = re.sub(r"[^a-z0-9]+", " ", value.lower())
    return " ".join(value.split())


def plain_text(value: str) -> str:
    """Remove the Markdown syntax used by the tutorials while preserving readable text."""
    value = re.sub(r"!\[([^]]*)\]\([^)]+\)", r"\1", value)
    value = re.sub(r"\[([^]]+)\]\([^)]+\)", r"\1", value)
    value = re.sub(r"<[^>]+>", "", value)
    value = re.sub(r"[*_`]", "", value)
    value = re.sub(r"^\s*>\s?", "", value)
    value = re.sub(r"\\([{}])", r"\1", value)
    return " ".join(value.split())


def shorten(value: str, max_words: int) -> str:
    words = plain_text(value).split()
    if len(words) <= max_words:
        return " ".join(words)
    return " ".join(words[:max_words]).rstrip(".,;:") + "…"


def shorten_code(value: str, max_characters: int) -> str:
    """Limit a code sample without applying prose-oriented Markdown cleanup."""
    if len(value) <= max_characters:
        return value
    return value[: max_characters - 1].rstrip() + "…"


@dataclass(frozen=True)
class MarkdownSection:
    heading: str
    level: int
    lines: tuple[str, ...]

    def paragraphs(self) -> list[str]:
        paragraphs: list[str] = []
        pending: list[str] = []
        in_fence = False
        for raw_line in self.lines:
            line = raw_line.strip()
            if line.startswith("```"):
                in_fence = not in_fence
                if pending:
                    paragraphs.append(plain_text(" ".join(pending)))
                    pending = []
                continue
            if in_fence or line.startswith("|") or re.match(r"^[-*+]\s+", line):
                if pending:
                    paragraphs.append(plain_text(" ".join(pending)))
                    pending = []
                continue
            if not line:
                if pending:
                    paragraphs.append(plain_text(" ".join(pending)))
                    pending = []
                continue
            if line.startswith("[") and "Tutorial home" in line:
                continue
            pending.append(line)
        if pending:
            paragraphs.append(plain_text(" ".join(pending)))
        return [paragraph for paragraph in paragraphs if paragraph]

    def bullets(self) -> list[str]:
        return [
            plain_text(match.group(1))
            for line in self.lines
            if (match := re.match(r"^\s*(?:[-*+]\s+|\d+\.\s+)(.+)$", line))
        ]

    def code_blocks(self) -> list[tuple[str, str]]:
        blocks: list[tuple[str, str]] = []
        language = ""
        lines: list[str] = []
        in_fence = False
        for raw_line in self.lines:
            if raw_line.strip().startswith("```"):
                if in_fence:
                    blocks.append((language, "\n".join(lines).strip()))
                    lines = []
                    language = ""
                    in_fence = False
                else:
                    language = raw_line.strip()[3:].strip().lower()
                    in_fence = True
                continue
            if in_fence:
                lines.append(raw_line)
        return blocks

    def tables(self) -> list[list[list[str]]]:
        tables: list[list[list[str]]] = []
        current: list[list[str]] = []
        for raw_line in (*self.lines, ""):
            line = raw_line.strip()
            if line.startswith("|") and line.endswith("|"):
                cells = [plain_text(cell.strip()) for cell in line.strip("|").split("|")]
                if all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells):
                    continue
                current.append(cells)
            elif current:
                if len(current) >= 2:
                    tables.append(current)
                current = []
        return tables


@dataclass(frozen=True)
class MarkdownDocument:
    path: Path
    title: str
    introduction: str
    sections: tuple[MarkdownSection, ...]

    @classmethod
    def read(cls, path: Path) -> MarkdownDocument:
        if not path.is_file():
            raise FileNotFoundError(f"Required tutorial file was not found: {path}")
        text = path.read_text(encoding="utf-8")
        sections: list[MarkdownSection] = []
        title = ""
        intro_lines: list[str] = []
        current_heading = ""
        current_level = 0
        current_lines: list[str] = []
        for line in text.splitlines():
            heading_match = re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
            if not heading_match:
                if current_heading:
                    current_lines.append(line)
                elif title:
                    intro_lines.append(line)
                continue
            if current_heading:
                sections.append(
                    MarkdownSection(current_heading, current_level, tuple(current_lines))
                )
            current_level = len(heading_match.group(1))
            current_heading = plain_text(heading_match.group(2))
            current_lines = []
            if current_level == 1 and not title:
                title = current_heading
                current_heading = ""
                current_level = 0
        if current_heading:
            sections.append(MarkdownSection(current_heading, current_level, tuple(current_lines)))
        if not title:
            raise ValueError(f"Tutorial has no level-one title: {path}")
        introduction = MarkdownSection("Introduction", 1, tuple(intro_lines)).paragraphs()
        return cls(path, title, introduction[0] if introduction else "", tuple(sections))

    def section(self, heading_fragment: str) -> MarkdownSection:
        target = normalize_heading(heading_fragment)
        for section in self.sections:
            if target in normalize_heading(section.heading):
                return section
        raise ValueError(f"{self.path.name} is missing required section: {heading_fragment}")


class LayoutGuard:
    """Reject coordinates that would place visible content outside the slide canvas."""

    @staticmethod
    def check(name: str, x: float, y: float, width: float, height: float) -> None:
        tolerance = 0.005
        if min(x, y, width, height) < 0:
            raise ValueError(f"{name} has a negative coordinate or dimension")
        if x + width > SLIDE_WIDTH + tolerance or y + height > SLIDE_HEIGHT + tolerance:
            raise ValueError(f"{name} exceeds the {SLIDE_WIDTH} x {SLIDE_HEIGHT} slide canvas")


def estimate_font_size(
    text: str,
    width: float,
    height: float,
    preferred: float,
    minimum: float = 10.5,
) -> float:
    """Estimate a safe font size from text length and the available bounding box."""
    size = preferred
    while size > minimum:
        chars_per_line = max(8, int(width * 72 / (size * 0.52)))
        lines = 0
        for paragraph in text.splitlines() or [text]:
            lines += max(1, math.ceil(max(1, len(paragraph)) / chars_per_line))
        capacity = max(1, int(height * 72 / (size * 1.22)))
        if lines <= capacity:
            return size
        size -= 0.5
    return minimum


def set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BACKGROUND)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float = 14,
    color: str = BODY,
    bold: bool = False,
    font: str = FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
    adaptive: bool = False,
    min_size: float = 10.5,
    hyperlink: str | None = None,
    name: str = "text",
):
    LayoutGuard.check(name, x, y, width, height)
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(
        estimate_font_size(text, width - margin * 2, height - margin * 2, size, min_size)
        if adaptive
        else size
    )
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    if hyperlink:
        shape.click_action.hyperlink.address = hyperlink
    return shape


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: str = CARD,
    line: str | None = None,
    radius_name: str = "shape",
):
    LayoutGuard.check(radius_name, x, y, width, height)
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_card(slide, x: float, y: float, width: float, height: float, *, fill: str = CARD):
    return add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        width,
        height,
        fill=fill,
        line="334155",
        radius_name="card",
    )


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    *,
    fill: str = TEAL_DARK,
    color: str = WHITE,
):
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        width,
        0.38,
        fill=fill,
        radius_name="pill",
    )
    return add_text(
        slide,
        text.upper(),
        x + 0.08,
        y + 0.04,
        width - 0.16,
        0.25,
        size=10,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        name="pill label",
    )


def add_header(slide, title: str, kicker: str, slide_number: int) -> None:
    add_pill(slide, kicker, 0.62, 0.48, min(2.5, 0.72 + len(kicker) * 0.085))
    add_text(
        slide,
        title,
        0.62,
        0.94,
        11.7,
        0.65,
        size=32,
        min_size=27,
        color=WHITE,
        bold=True,
        adaptive=True,
        name="slide title",
    )
    add_text(
        slide,
        f"{slide_number:02d}",
        12.17,
        0.52,
        0.55,
        0.3,
        size=11,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
        name="slide number",
    )


def add_footer(slide, source: str) -> None:
    add_text(
        slide,
        source,
        0.62,
        7.14,
        12.05,
        0.2,
        size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        name="source footer",
    )


def add_card_text(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    title: str,
    body: str,
    *,
    accent: str = TEAL,
    body_size: float = 14,
) -> None:
    add_card(slide, x, y, width, height)
    add_text(
        slide,
        title,
        x + 0.28,
        y + 0.24,
        width - 0.56,
        0.4,
        size=18,
        color=accent,
        bold=True,
        adaptive=True,
        min_size=15,
        name=f"{title} heading",
    )
    add_text(
        slide,
        body,
        x + 0.28,
        y + 0.82,
        width - 0.56,
        height - 1.08,
        size=body_size,
        color=BODY,
        adaptive=True,
        min_size=11,
        name=f"{title} body",
    )


def add_code_box(slide, label: str, code: str, x: float, y: float, width: float, height: float):
    add_card(slide, x, y, width, height, fill="111827")
    add_text(
        slide,
        label.upper(),
        x + 0.24,
        y + 0.17,
        width - 0.48,
        0.25,
        size=9.5,
        color=TEAL,
        bold=True,
        name=f"{label} code label",
    )
    add_text(
        slide,
        code,
        x + 0.24,
        y + 0.54,
        width - 0.48,
        height - 0.72,
        size=12.5,
        min_size=9.5,
        color=BODY,
        font=CODE_FONT,
        adaptive=True,
        name=f"{label} code",
    )


def table_cell(table: list[list[str]], row_name: str, column_name: str) -> str:
    headers = [normalize_heading(cell) for cell in table[0]]
    column = normalize_heading(column_name)
    try:
        column_index = next(i for i, header in enumerate(headers) if column in header)
    except StopIteration as error:
        raise ValueError(f"Markdown table is missing column: {column_name}") from error
    row_target = normalize_heading(row_name)
    for row in table[1:]:
        if row and row_target in normalize_heading(row[0]) and column_index < len(row):
            return row[column_index]
    raise ValueError(f"Markdown table is missing row: {row_name}")


def first_paragraph(section: MarkdownSection, *, words: int = 34) -> str:
    paragraphs = section.paragraphs()
    if not paragraphs:
        raise ValueError(f"Section has no usable paragraph: {section.heading}")
    return shorten(paragraphs[0], words)


def find_code(section: MarkdownSection, language: str | None = None) -> str:
    for block_language, code in section.code_blocks():
        if language is None or language in block_language:
            return code
    raise ValueError(f"Section has no required code block: {section.heading}")


def add_forecast_trace(slide, x: float, y: float, width: float, height: float) -> None:
    """Draw a native vector history-to-forecast motif."""
    values = [0.66, 0.49, 0.55, 0.34, 0.44, 0.25, 0.30, 0.15]
    points: list[tuple[float, float]] = []
    for index, value in enumerate(values):
        px = x + index * width / (len(values) - 1)
        py = y + value * height
        points.append((px, py))
    for index in range(len(points) - 1):
        x1, y1 = points[index]
        x2, y2 = points[index + 1]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        connector.line.color.rgb = rgb(TEAL if index >= 4 else MUTED)
        connector.line.width = Pt(2.5)
    for index, (px, py) in enumerate(points):
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            px - 0.07,
            py - 0.07,
            0.14,
            0.14,
            fill=TEAL if index >= 5 else "64748B",
            radius_name="forecast node",
        )


def create_title_slide(prs: Presentation, intro: MarkdownDocument) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_pill(slide, "Zero to Master Course", 0.72, 0.72, 2.28)
    add_text(
        slide,
        "TimesFM\nForecasting Studio",
        0.72,
        1.48,
        7.15,
        1.55,
        size=38,
        min_size=34,
        color=WHITE,
        bold=True,
        adaptive=True,
        name="deck title",
    )
    tagline = intro.introduction or "Local-first, zero-shot forecasting with Google TimesFM."
    add_text(
        slide,
        shorten(tagline, 24),
        0.75,
        3.25,
        6.65,
        0.72,
        size=18,
        min_size=15,
        color=BODY,
        adaptive=True,
        name="deck tagline",
    )
    add_pill(slide, "Local-first", 0.75, 4.35, 1.32, fill="334155")
    add_pill(slide, "Zero-shot", 2.22, 4.35, 1.28, fill="334155")
    add_pill(slide, "Univariate", 3.65, 4.35, 1.28, fill="334155")
    add_card(slide, 8.12, 1.14, 4.42, 4.72, fill=CARD_ALT)
    add_text(
        slide,
        "history",
        8.52,
        1.58,
        1.2,
        0.3,
        size=11,
        color=MUTED,
        bold=True,
        name="history label",
    )
    add_text(
        slide,
        "forecast",
        10.78,
        1.58,
        1.18,
        0.3,
        size=11,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.RIGHT,
        name="forecast label",
    )
    add_forecast_trace(slide, 8.62, 2.08, 3.38, 2.46)
    add_text(
        slide,
        "CONTEXT  →  HORIZON",
        8.62,
        4.92,
        3.38,
        0.34,
        size=12,
        color=BODY,
        bold=True,
        align=PP_ALIGN.CENTER,
        name="context horizon label",
    )
    add_text(
        slide,
        "Google TimesFM 2.5 · Streamlit · Plotly",
        0.72,
        6.78,
        7.5,
        0.28,
        size=11,
        color=MUTED,
        name="technology caption",
    )
    add_text(
        slide,
        "01 / 08",
        11.63,
        6.78,
        0.9,
        0.28,
        size=11,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
        name="title slide number",
    )


def create_paradigm_slide(prs: Presentation, intro: MarkdownDocument) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "The paradigm shift", "Model thinking", 2)
    comparison = intro.section("TimesFM versus established approaches").tables()[0]
    traditional = (
        f"ARIMA · {table_cell(comparison, 'Core mechanism', 'ARIMA')}\n\n"
        f"Prophet · {table_cell(comparison, 'Core mechanism', 'Prophet')}\n\n"
        f"Fitting · {table_cell(comparison, 'Typical fitting', 'ARIMA')}"
    )
    timesfm = (
        f"Engine · {table_cell(comparison, 'Core mechanism', 'TimesFM')}\n\n"
        f"Reuse · {table_cell(comparison, 'Typical fitting', 'TimesFM')}\n\n"
        f"Uncertainty · {table_cell(comparison, 'Uncertainty', 'TimesFM')}"
    )
    add_card_text(slide, 0.65, 1.86, 5.35, 4.62, "Fit per series", traditional, accent=BLUE)
    add_card_text(slide, 7.33, 1.86, 5.35, 4.62, "Pretrain once · forecast zero-shot", timesfm)
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        6.15,
        3.31,
        1.05,
        1.05,
        fill=TEAL_DARK,
        line=TEAL,
        radius_name="shift icon",
    )
    add_text(
        slide,
        "→",
        6.33,
        3.48,
        0.68,
        0.44,
        size=24,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        name="shift arrow",
    )
    add_pill(slide, "No target-domain fitting", 5.35, 5.78, 2.63, fill="334155")
    add_footer(slide, "Source: 01_timesfm_intro.md · sections 3-4")


def create_architecture_slide(
    prs: Presentation, data_doc: MarkdownDocument, forecast_doc: MarkdownDocument
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Core application architecture", "Data to decision", 3)
    descriptions = [
        first_paragraph(data_doc.section("Ingestion architecture"), words=19),
        "Select the date and target, parse timestamps, sort, and validate an exact regular grid.",
        first_paragraph(forecast_doc.section("End-to-end prediction path"), words=20),
        "Attach future timestamps and render the q50 path with the q10-q90 uncertainty band.",
    ]
    labels = ["Data ingestion", "Datetime parsing", "TimesFM inference", "Plotly output"]
    symbols = ["01", "02", "03", "04"]
    x_positions = [0.64, 3.82, 7.0, 10.18]
    for index, (x, label, description, symbol) in enumerate(
        zip(x_positions, labels, descriptions, symbols, strict=True)
    ):
        add_card(slide, x, 2.15, 2.5, 3.66, fill=CARD if index % 2 == 0 else CARD_ALT)
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            x + 0.24,
            2.42,
            0.54,
            0.54,
            fill=TEAL_DARK,
            radius_name=f"architecture node {index + 1}",
        )
        add_text(
            slide,
            symbol,
            x + 0.24,
            2.56,
            0.54,
            0.18,
            size=10,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"architecture number {index + 1}",
        )
        add_text(
            slide,
            label,
            x + 0.24,
            3.25,
            2.02,
            0.72,
            size=19,
            min_size=16,
            color=WHITE,
            bold=True,
            adaptive=True,
            name=f"architecture label {index + 1}",
        )
        add_text(
            slide,
            shorten(description, 25),
            x + 0.24,
            4.2,
            2.02,
            1.22,
            size=13,
            min_size=10.5,
            color=BODY,
            adaptive=True,
            name=f"architecture description {index + 1}",
        )
        if index < 3:
            add_shape(
                slide,
                MSO_SHAPE.CHEVRON,
                x + 2.58,
                3.54,
                0.42,
                0.68,
                fill=TEAL,
                radius_name=f"architecture chevron {index + 1}",
            )
    add_pill(slide, "One validated univariate series", 5.03, 6.15, 3.24, fill="334155")
    add_footer(slide, "Sources: 03_data_engineering.md · 04_forecasting_mastery.md")


def create_ingestion_slide(prs: Presentation, data_doc: MarkdownDocument) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Ingestion and APIs", "Bring your data", 4)
    format_table = data_doc.section("Supported file formats").tables()[0]
    formats = []
    for row in format_table[1:4]:
        formats.append((row[0], row[2], row[3]))
    x_positions = [0.66, 2.76, 4.86]
    for x, (name, strength, warning) in zip(x_positions, formats, strict=True):
        add_card(slide, x, 2.0, 1.78, 3.75)
        add_text(
            slide,
            name,
            x + 0.2,
            2.28,
            1.38,
            0.42,
            size=20,
            color=TEAL,
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"{name} format",
        )
        add_text(
            slide,
            shorten(strength, 14),
            x + 0.2,
            3.02,
            1.38,
            0.9,
            size=13,
            min_size=11,
            color=BODY,
            adaptive=True,
            align=PP_ALIGN.CENTER,
            name=f"{name} strength",
        )
        add_text(
            slide,
            shorten(warning, 13),
            x + 0.2,
            4.48,
            1.38,
            0.72,
            size=11.5,
            min_size=10,
            color=MUTED,
            adaptive=True,
            align=PP_ALIGN.CENTER,
            name=f"{name} warning",
        )
    add_card_text(
        slide,
        7.05,
        2.0,
        2.62,
        3.75,
        "Kaggle",
        first_paragraph(data_doc.section("Retrieving Kaggle datasets"), words=30),
        accent=AMBER,
        body_size=13,
    )
    add_card_text(
        slide,
        10.05,
        2.0,
        2.62,
        3.75,
        "Hugging Face",
        first_paragraph(data_doc.section("Retrieving Hugging Face datasets"), words=30),
        body_size=13,
    )
    add_pill(slide, "Local uploads", 0.66, 6.12, 1.38, fill="334155")
    add_pill(slide, "Public HTTP/S URLs", 2.22, 6.12, 1.92, fill="334155")
    add_pill(slide, "Credential-safe providers", 4.32, 6.12, 2.38, fill="334155")
    add_footer(slide, "Source: 03_data_engineering.md · sections 2, 4, 7, 8")


def create_studio_slide(prs: Presentation, forecast_doc: MarkdownDocument) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Interactive studio features", "Tune and test", 5)
    context_table = forecast_doc.section("Choosing context by cycles").tables()[0]
    horizon_table = forecast_doc.section("Forecast horizon").tables()[0]
    add_card(slide, 0.64, 1.9, 7.22, 4.9)
    add_text(
        slide,
        "Context  C",
        0.98,
        2.22,
        2.2,
        0.42,
        size=19,
        color=TEAL,
        bold=True,
        name="context label",
    )
    add_text(
        slide,
        "How much recent history the model receives",
        0.98,
        2.76,
        3.05,
        0.58,
        size=13,
        color=BODY,
        adaptive=True,
        name="context definition",
    )
    add_text(
        slide,
        table_cell(context_table, "Hourly", "Candidate context"),
        4.8,
        2.16,
        2.2,
        0.72,
        size=38,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.RIGHT,
        name="context metric",
    )
    add_text(
        slide,
        table_cell(context_table, "Hourly", "Patterns represented"),
        4.12,
        2.94,
        2.88,
        0.38,
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        name="context metric caption",
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.98, 3.62, 6.02, 0.03, fill="334155")
    add_text(
        slide,
        "Horizon  H",
        0.98,
        4.04,
        2.2,
        0.42,
        size=19,
        color=BLUE,
        bold=True,
        name="horizon label",
    )
    add_text(
        slide,
        "How many future periods the model returns",
        0.98,
        4.58,
        3.05,
        0.58,
        size=13,
        color=BODY,
        adaptive=True,
        name="horizon definition",
    )
    add_text(
        slide,
        table_cell(horizon_table, "Hourly", "Horizon"),
        4.8,
        3.98,
        2.2,
        0.72,
        size=38,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.RIGHT,
        name="horizon metric",
    )
    add_text(
        slide,
        table_cell(horizon_table, "Hourly", "Calendar interpretation"),
        4.12,
        4.76,
        2.88,
        0.38,
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        name="horizon metric caption",
    )
    add_pill(slide, "Interactive controls", 0.98, 5.62, 1.82, fill=TEAL_DARK)
    manual = forecast_doc.section("Manual simulator")
    manual_code = find_code(manual)
    add_code_box(slide, "Manual simulator", manual_code, 8.2, 1.9, 4.48, 2.48)
    parser_steps = manual.bullets()[:3]
    add_card_text(
        slide,
        8.2,
        4.72,
        4.48,
        2.08,
        "Fast-path forecast",
        "\n".join(f"{index + 1}. {shorten(step, 16)}" for index, step in enumerate(parser_steps)),
        body_size=13,
    )
    add_footer(slide, "Source: 04_forecasting_mastery.md · sections 2, 3, 8, 9")


def repository_name(github_url: str) -> str:
    path = urlparse(github_url).path.rstrip("/")
    name = path.rsplit("/", 1)[-1]
    return name.removesuffix(".git") or "repository"


def create_setup_slide(prs: Presentation, install_doc: MarkdownDocument, github_url: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Local setup and installation", "Three commands", 6)
    repo_name = repository_name(github_url)
    clone_code = f"git clone {github_url}\nSet-Location {repo_name}"
    uv_code = find_code(install_doc.section("Recommended setup with uv"), "powershell")
    uv_lines = uv_code.splitlines()[:2]
    run_code = "uv run streamlit run app.py"
    steps = [
        ("01", "Clone", clone_code),
        ("02", "Create environment", "\n".join(uv_lines)),
        ("03", "Run", run_code),
    ]
    for index, (number, label, code) in enumerate(steps):
        x = 0.66 + index * 4.18
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            x,
            1.9,
            0.56,
            0.56,
            fill=TEAL_DARK,
            radius_name=f"setup step {number}",
        )
        add_text(
            slide,
            number,
            x,
            2.05,
            0.56,
            0.18,
            size=10,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"setup step number {number}",
        )
        add_text(
            slide,
            label,
            x + 0.76,
            1.98,
            2.88,
            0.4,
            size=19,
            color=WHITE,
            bold=True,
            name=f"setup label {number}",
        )
        add_code_box(slide, label, code, x, 2.72, 3.64, 2.15)
        if index < 2:
            add_shape(
                slide,
                MSO_SHAPE.CHEVRON,
                x + 3.76,
                3.42,
                0.32,
                0.55,
                fill=TEAL,
                radius_name=f"setup chevron {index + 1}",
            )
    cuda_code = find_code(install_doc.section("Verify PyTorch"), "powershell")
    cuda_code = shorten_code(cuda_code.replace("uv run python -c ", ""), 240)
    add_card(slide, 0.66, 5.35, 12.02, 1.36, fill="172033")
    add_text(
        slide,
        "GPU CHECK",
        0.96,
        5.67,
        1.2,
        0.3,
        size=11,
        color=AMBER,
        bold=True,
        name="gpu check label",
    )
    add_text(
        slide,
        cuda_code,
        2.25,
        5.58,
        9.96,
        0.52,
        size=11.5,
        min_size=9.5,
        color=BODY,
        font=CODE_FONT,
        adaptive=True,
        name="gpu check command",
    )
    add_footer(slide, "Source: 02_local_installation.md · sections 3, 4, 7")


def create_pitfalls_slide(
    prs: Presentation, data_doc: MarkdownDocument, forecast_doc: MarkdownDocument
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Performance and pitfalls", "Operate safely", 7)
    add_card(slide, 0.65, 1.9, 4.03, 4.9)
    add_text(
        slide,
        "16,384",
        0.95,
        2.28,
        3.43,
        0.82,
        size=43,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        name="compile limit metric",
    )
    add_text(
        slide,
        "rounded context + horizon limit",
        0.95,
        3.16,
        3.43,
        0.4,
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        name="compile limit caption",
    )
    add_text(
        slide,
        "Ccompile = 32⌈C / 32⌉\nHcompile = 128⌈H / 128⌉",
        1.0,
        4.02,
        3.33,
        0.9,
        size=18,
        min_size=15,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
        adaptive=True,
        name="compile formulas",
    )
    add_pill(slide, "Patch boundaries matter", 1.38, 5.57, 2.58, fill="334155")
    missing_table = data_doc.section("Interpolation mathematics").tables()[0]
    trailing = table_cell(missing_table, "Trailing target gap", "App behavior")
    warnings = [
        (
            "Missing or uninitialized data",
            f"{trailing}. Internal gaps are interpolated only when both anchors exist.",
            AMBER,
        ),
        (
            "GPU out of memory",
            first_paragraph(forecast_doc.section("Performance and OOM control"), words=31),
            ROSE,
        ),
        (
            "Frequency mismatch",
            "Frequency validates timestamps and creates the future index; "
            "TimesFM 2.5 does not receive a frequency token.",
            BLUE,
        ),
    ]
    for index, (title, body, accent) in enumerate(warnings):
        y = 1.9 + index * 1.66
        add_card(slide, 5.05, y, 7.63, 1.38, fill=CARD_ALT)
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            5.34,
            y + 0.36,
            0.56,
            0.56,
            fill=accent,
            radius_name=f"warning icon {index + 1}",
        )
        add_text(
            slide,
            "!",
            5.34,
            y + 0.49,
            0.56,
            0.23,
            size=14,
            color=BACKGROUND,
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"warning mark {index + 1}",
        )
        add_text(
            slide,
            title,
            6.18,
            y + 0.2,
            5.98,
            0.35,
            size=17,
            color=accent,
            bold=True,
            name=f"warning heading {index + 1}",
        )
        add_text(
            slide,
            shorten(body, 29),
            6.18,
            y + 0.65,
            5.98,
            0.48,
            size=12.5,
            min_size=10.5,
            color=BODY,
            adaptive=True,
            name=f"warning body {index + 1}",
        )
    add_footer(slide, "Sources: 03_data_engineering.md · 04_forecasting_mastery.md")


def create_cta_slide(prs: Presentation, github_url: str, docs_url: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_pill(slide, "Your next forecast", 0.72, 0.68, 1.82)
    add_text(
        slide,
        "Build the baseline.\nChallenge the forecast.",
        0.72,
        1.42,
        7.45,
        1.4,
        size=35,
        min_size=31,
        color=WHITE,
        bold=True,
        adaptive=True,
        name="cta title",
    )
    add_text(
        slide,
        "Move from a clean time series to a reproducible zero-shot experiment.",
        0.75,
        3.05,
        6.75,
        0.6,
        size=17,
        color=BODY,
        adaptive=True,
        name="cta subtitle",
    )
    next_steps = [
        ("01", "Clone", "Get the local-first studio running."),
        ("02", "Forecast", "Load a regular series and tune C versus H."),
        ("03", "Validate", "Compare against naive and statistical baselines."),
    ]
    for index, (number, title, body) in enumerate(next_steps):
        x = 0.74 + index * 2.45
        add_card(slide, x, 4.18, 2.14, 1.63, fill=CARD_ALT)
        add_text(
            slide,
            number,
            x + 0.22,
            4.42,
            0.4,
            0.3,
            size=11,
            color=TEAL,
            bold=True,
            name=f"cta step number {number}",
        )
        add_text(
            slide,
            title,
            x + 0.65,
            4.38,
            1.2,
            0.35,
            size=17,
            color=WHITE,
            bold=True,
            name=f"cta step title {number}",
        )
        add_text(
            slide,
            body,
            x + 0.22,
            4.93,
            1.7,
            0.58,
            size=11.5,
            min_size=10,
            color=BODY,
            adaptive=True,
            name=f"cta step body {number}",
        )
    add_card(slide, 8.55, 1.38, 4.05, 4.9, fill=CARD_ALT)
    add_text(
        slide,
        "PROJECT LINKS",
        8.92,
        1.82,
        3.31,
        0.32,
        size=11,
        color=MUTED,
        bold=True,
        name="project links label",
    )
    add_text(
        slide,
        "GitHub repository ↗",
        8.92,
        2.48,
        3.15,
        0.46,
        size=19,
        color=TEAL,
        bold=True,
        hyperlink=github_url,
        name="github link",
    )
    add_text(
        slide,
        shorten(github_url, 8),
        8.92,
        3.02,
        3.15,
        0.48,
        size=10.5,
        min_size=8.5,
        color=MUTED,
        adaptive=True,
        hyperlink=github_url,
        name="github url",
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.92, 3.68, 3.15, 0.03, fill="334155")
    add_text(
        slide,
        "Course documentation ↗",
        8.92,
        4.12,
        3.15,
        0.46,
        size=19,
        min_size=16,
        color=BLUE,
        bold=True,
        adaptive=True,
        hyperlink=docs_url,
        name="documentation link",
    )
    add_text(
        slide,
        shorten(docs_url, 8),
        8.92,
        4.67,
        3.15,
        0.48,
        size=10.5,
        min_size=8.5,
        color=MUTED,
        adaptive=True,
        hyperlink=docs_url,
        name="documentation url",
    )
    add_pill(slide, "Zero-shot ≠ zero-validation", 9.02, 5.52, 2.95, fill=TEAL_DARK)
    add_text(
        slide,
        "08 / 08",
        11.63,
        6.78,
        0.9,
        0.28,
        size=11,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.RIGHT,
        name="cta slide number",
    )


def validate_http_url(value: str, option_name: str) -> str:
    parsed = urlparse(value)
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        raise argparse.ArgumentTypeError(f"{option_name} must be an absolute HTTP/S URL")
    if parsed.username or parsed.password:
        raise argparse.ArgumentTypeError(f"{option_name} must not contain credentials")
    return value


def load_documents(source_dir: Path) -> tuple[MarkdownDocument, ...]:
    names = (
        "01_timesfm_intro.md",
        "02_local_installation.md",
        "03_data_engineering.md",
        "04_forecasting_mastery.md",
    )
    return tuple(MarkdownDocument.read(source_dir / name) for name in names)


def build_presentation(
    documents: tuple[MarkdownDocument, ...], github_url: str, docs_url: str
) -> Presentation:
    if len(documents) != 4:
        raise ValueError("Exactly four tutorial documents are required")
    intro, install_doc, data_doc, forecast_doc = documents
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    prs.core_properties.title = "TimesFM Forecasting Studio"
    prs.core_properties.subject = "Zero to Master Course"
    prs.core_properties.author = "TimesFM Forecasting Studio"
    prs.core_properties.keywords = "TimesFM, forecasting, Streamlit, time series"
    create_title_slide(prs, intro)
    create_paradigm_slide(prs, intro)
    create_architecture_slide(prs, data_doc, forecast_doc)
    create_ingestion_slide(prs, data_doc)
    create_studio_slide(prs, forecast_doc)
    create_setup_slide(prs, install_doc, github_url)
    create_pitfalls_slide(prs, data_doc, forecast_doc)
    create_cta_slide(prs, github_url, docs_url)
    if len(prs.slides) != 8:
        raise RuntimeError(f"Expected 8 slides, generated {len(prs.slides)}")
    return prs


def save_atomically(prs: Presentation, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    temporary_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            prefix=f".{output_path.stem}-",
            suffix=".pptx",
            dir=output_path.parent,
            delete=False,
        ) as temporary_file:
            temporary_path = Path(temporary_file.name)
        prs.save(temporary_path)
        reopened = Presentation(temporary_path)
        if len(reopened.slides) != 8:
            raise RuntimeError("Saved presentation did not reopen with exactly 8 slides")
        os.replace(temporary_path, output_path)
        temporary_path = None
    finally:
        if temporary_path is not None and temporary_path.exists():
            temporary_path.unlink()


def parse_arguments() -> argparse.Namespace:
    base_dir = Path(__file__).resolve().parent
    parser = argparse.ArgumentParser(
        description="Generate an 8-slide TimesFM Masterclass presentation from tutorial Markdown."
    )
    parser.add_argument(
        "--github-url",
        default="https://github.com/pypi-ahmad/timesfm-forecasting-studio",
        type=lambda value: validate_http_url(value, "--github-url"),
        help="Published project GitHub URL used on the setup and CTA slides.",
    )
    parser.add_argument(
        "--docs-url",
        default="https://github.com/pypi-ahmad/timesfm-forecasting-studio/tree/main/docs/tutorial",
        type=lambda value: validate_http_url(value, "--docs-url"),
        help="Published course documentation URL used on the CTA slide.",
    )
    parser.add_argument(
        "--source-dir",
        type=Path,
        default=base_dir / "docs" / "tutorial",
        help="Directory containing the four tutorial Markdown files.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=base_dir / "TimesFM_Masterclass.pptx",
        help="Destination .pptx path.",
    )
    arguments = parser.parse_args()
    if arguments.output.suffix.lower() != ".pptx":
        parser.error("--output must end with .pptx")
    return arguments


def main() -> None:
    arguments = parse_arguments()
    source_dir = arguments.source_dir.expanduser().resolve()
    output_path = arguments.output.expanduser().resolve()
    documents = load_documents(source_dir)
    presentation = build_presentation(documents, arguments.github_url, arguments.docs_url)
    save_atomically(presentation, output_path)
    print(f"Successfully created {output_path.name}")


if __name__ == "__main__":
    main()
